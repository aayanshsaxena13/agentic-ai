# Import statements here...
from pptx import Presentation
from pptx.util import Pt, Inches
import os
import json

# Create the required variables...
input_path = os.path.join(os.path.dirname(__file__), "..", "input", "powerpoint_input.json")
output_path = os.path.join(os.path.dirname(__file__), "..", "output", "powerpoint.pptx")

# Create the module function here...
def createTitleSlide():
    with open(input_path, "r") as f:
        file = json.loads(f.read())
    
    prs = Presentation()
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout=slide_layout)

    title = slide.shapes.title
    title.text = file["title"]

    subtitle = slide.placeholders[1]
    subtitle.text = file["subtitle"]

    prs.save(output_path)
    return "Generated powerpoint title slide using Agentic AI..."

def createContentSlide():
    prs = Presentation(output_path)

    with open(input_path, "r") as f:
        file = json.loads(f.read())

    bullet_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_layout)

    title = slide.shapes.title
    title.text = file["title"]

    content = slide.placeholders[1]
    content.text = file["content"]

    prs.save(output_path)

    return "Generated powerpoint content slide by Agentic AI..."